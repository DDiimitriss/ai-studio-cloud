import os
import re
from pptx import Presentation
from pptx.util import Inches

def run(args):
    html_content = args.get("html", "")
    title = args.get("title", "Presentation")
    if not html_content:
        return "Error: No HTML content provided."

    try:
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # Title and Content

        # Very simple extraction: each <h2> becomes a slide title, following text as content
        # This is basic; for production use a proper HTML parser like BeautifulSoup.
        # Split by <h2> tags
        parts = re.split(r'<h2[^>]*>(.*?)</h2>', html_content, flags=re.DOTALL)
        # parts[0] is text before first h2, then alternating: title, content, title, content...
        for i in range(1, len(parts), 2):
            title_text = re.sub(r'<[^>]+>', '', parts[i]).strip()
            content_text = parts[i+1] if i+1 < len(parts) else ""
            content_text = re.sub(r'<[^>]+>', '', content_text).strip()
            if not title_text:
                continue
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_text[:100]
            # Content placeholder may be index 1
            if slide.placeholders and len(slide.placeholders) > 1:
                slide.placeholders[1].text = content_text[:500] if content_text else "(no content)"
        desktop = os.path.join(os.environ["USERPROFILE"], "Desktop")
        safe_title = re.sub(r'[\\/*?:"<>|]', "", title)
        filename = f"{safe_title}.pptx"
        save_path = os.path.join(desktop, filename)
        prs.save(save_path)
        return f"PowerPoint saved to {save_path}"
    except Exception as e:
        return f"PPT creation error: {e}"

def get_info():
    return {
        "name": "Doc Presentation",
        "description": "Converts HTML to PowerPoint (.pptx). Arguments: 'html', 'title'. Use when user wants a presentation from HTML outline."
    }